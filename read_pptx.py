import pptx
import os

pptx_file = 'Presentación del PI M5 (1).pptx'
output_file = 'pptx_content_utf8.txt'

if not os.path.exists(pptx_file):
    print(f"File not found: {pptx_file}")
    exit(1)

try:
    prs = pptx.Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write("\n".join(text))
    print(f"Successfully wrote to {output_file}")
except Exception as e:
    print(f"Error reading pptx: {e}")
